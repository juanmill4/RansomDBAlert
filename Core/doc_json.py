import re
import hashlib
import json
import os
from concurrent.futures import ThreadPoolExecutor

# Para leer archivos DOCX
import docx
# Para leer archivos PPTX
from pptx import Presentation

# Expresión regular para extraer emails (RFC 2822)
email_regex = re.compile(
    r"([^\x00-\x20\x22\x28\x29\x2c\x2e\x3a-\x3c\x3e\x40\x5b-\x5d\x7f-\xff\x7c]+|\x22([^\x0d\x22\x5c\x80-\xff]"
    r"|\x5c[\x00-\x7f])*\x22)(\x2e([^\x00-\x20\x22\x28\x29\x2c\x2e\x3a-\x3c\x3e\x40\x5b-\x5d\x7f-\xff\x7c]+"
    r"|\x22([^\x0d\x22\x5c\x80-\xff]|\x5c[\x00-\x7f])*\x22))*\x40([^\x00-\x20\x22\x28\x29\x2c\x2e\x3a-\x3c"
    r"\x3e\x40\x5b-\x5d\x7f-\xff\x7c]+|\x5b([^\x0d\x5b-\x5d\x80-\xff]|\x5c[\x00-\x7f])*\x5d)(\x2e([^\x00-\x20"
    r"\x22\x28\x29\x2c\x2e\x3a-\x3c\x3e\x40\x5b-\x5d\x7f-\xff\x7c]+|\x5b([^\x0d\x5b-\x5d\x80-\xff]|"
    r"\x5c[\x00-\x7f])*\x5d))"
)

def generar_id(email, contexto):
    """Genera un hash SHA-1 único basado en email + contexto."""
    return hashlib.sha1(f"{email}{contexto}".encode()).hexdigest()

def extraer_emails_de_texto(texto_completo, fuente):
    """
    Función genérica para buscar emails en un texto y devolver un dict
    con email, contexto y ID.
    - texto_completo (str): texto donde buscar.
    - fuente (str): puede ser "DOCX Archive", "PPTX Archive", etc.
    """
    resultados = {}
    for match in email_regex.finditer(texto_completo):
        email = match.group(0)
        start = max(0, match.start() - 500)
        end = min(len(texto_completo), match.end() + 500)
        contexto = texto_completo[start:end].replace("\n", " ")
        # Generar ID único
        id_hash = generar_id(email, contexto)
        resultados[email] = {
            "email_context": contexto,
            "ID": id_hash,
            "FROM": fuente
        }
    return resultados

def extraer_emails_con_contexto_docx(docx_file, directorio, directorio_json):
    """
    Extrae emails con contexto de un archivo DOCX.
    """
    ruta_docx = os.path.join(directorio, docx_file)
    print(f"Procesando DOCX: {ruta_docx}")
    
    # Abrimos el DOCX
    doc = docx.Document(ruta_docx)
    
    # Convertimos todo el contenido en un gran string
    texto_completo = []
    
    # Recorremos párrafos
    for para in doc.paragraphs:
        texto_completo.append(para.text)
    
    # Recorremos tablas con try-except
    for table in doc.tables:
        try:
            for row in table.rows:
                for cell in row.cells:
                    texto_completo.append(cell.text)
        except ValueError as e:
            # Si quisieras "loggear" algo, podrías hacerlo aquí
            print(f"Advertencia: se ignoró ValueError en {ruta_docx}: {e}")
            # Continúas sin interrumpir el programa
            continue
    
    texto_unido = "\n".join(texto_completo)
    
    # Buscamos emails
    resultados = extraer_emails_de_texto(texto_unido, "DOCX Archive")
    
    # Guardamos si hay resultados
    if resultados:
        nombre_json = os.path.splitext(os.path.basename(ruta_docx))[0] + ".json"
        ruta_json = os.path.join(directorio_json, nombre_json)
        with open(ruta_json, "w", encoding="utf-8") as f:
            json.dump(resultados, f, indent=4, ensure_ascii=False)
        print(f"Procesado: {ruta_docx} → {ruta_json} ({len(resultados)} emails extraídos)")
        os.remove(ruta_docx)
    else:
        print(f"No se encontraron emails en {ruta_docx}")
        os.remove(ruta_docx)


def extraer_emails_con_contexto_pptx(pptx_file, directorio, directorio_json):
    """
    Extrae emails con contexto de un archivo PPTX.
    """
    ruta_pptx = os.path.join(directorio, pptx_file)
    print(f"Procesando PPTX: {ruta_pptx}")
    
    prs = Presentation(ruta_pptx)
    texto_completo = []
    
    # Recorremos slides y shapes para extraer texto
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto_completo.append(shape.text)
    
    # Unimos todo
    texto_unido = "\n".join(texto_completo)
    
    # Buscamos emails
    resultados = extraer_emails_de_texto(texto_unido, "PPTX Archive")
    
    if resultados:
        nombre_json = os.path.splitext(os.path.basename(ruta_pptx))[0] + ".json"
        ruta_json = os.path.join(directorio_json, nombre_json)
        with open(ruta_json, "w", encoding="utf-8") as f:
            json.dump(resultados, f, indent=4, ensure_ascii=False)
        print(f"Procesado: {ruta_pptx} → {ruta_json} ({len(resultados)} emails extraídos)")
    else:
        print(f"No se encontraron emails en {ruta_pptx}")

    os.remove(ruta_pptx)

def docs_to_json(directorio):
    """
    Procesa los archivos DOCX y PPTX en un directorio y
    guarda los JSON en subcarpetas separadas.
    """
    # Directorios de salida
    directorio_docx_json = os.path.join(directorio, "DOCX_to_json")
    directorio_pptx_json = os.path.join(directorio, "PPTX_to_json")
    
    os.makedirs(directorio_docx_json, exist_ok=True)
    os.makedirs(directorio_pptx_json, exist_ok=True)
    
    # Listamos archivos
    archivos = [f for f in os.listdir(directorio) if os.path.isfile(os.path.join(directorio, f))]
    
    docx_files = [f for f in archivos if f.lower().endswith(".docx")]
    pptx_files = [f for f in archivos if f.lower().endswith(".pptx") or f.lower().endswith(".ppt")]
    # Si quieres excluir .ppt “antiguo” y procesar sólo .pptx, ajusta lo anterior

    if not docx_files and not pptx_files:
        print("No se encontraron archivos DOCX ni PPTX en el directorio.")
        return

    print(f"Se encontraron {len(docx_files)} archivos DOCX y {len(pptx_files)} archivos PPT/PPTX. Procesando...\n")

    # Procesamiento en paralelo
    with ThreadPoolExecutor() as executor:
        # Procesar DOCX
        futures_docx = [
            executor.submit(extraer_emails_con_contexto_docx, f, directorio, directorio_docx_json)
            for f in docx_files
        ]
        # Procesar PPTX
        futures_pptx = [
            executor.submit(extraer_emails_con_contexto_pptx, f, directorio, directorio_pptx_json)
            for f in pptx_files
        ]
        # Esperamos a que terminen
        for future in futures_docx + futures_pptx:
            future.result()

    print("Todos los archivos DOCX/PPTX han sido procesados. Los JSON están en sus carpetas respectivas.")

# ==== EJEMPLO DE USO ====
if __name__ == "__main__":
    ruta_archivos = "Dir_documentos/digitalizados"  # Cambia esta ruta a tu carpeta de .docx/.pptx
    docs_to_json(ruta_archivos)
