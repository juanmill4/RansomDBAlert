import os
import subprocess
import shutil
from concurrent.futures import ThreadPoolExecutor
from docx import Document
from pptx import Presentation

def convert_doc_to_docx(doc_file):
    """
    Convierte un archivo .doc a .docx utilizando LibreOffice en modo headless.
    Retorna la ruta del nuevo archivo .docx o None si falla.
    """
    try:
        parent_dir = os.path.dirname(doc_file)
        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to",
            "docx",
            doc_file
        ], cwd=parent_dir, check=True, timeout=6)  
        
        base, _ = os.path.splitext(doc_file)
        new_file = base + ".docx"
        if os.path.exists(new_file):
            return new_file
            os.remove(doc_file)
        else:
            return None
    
    except subprocess.TimeoutExpired:
        print(f"Tiempo de espera excedido al convertir {doc_file}")
        return None
    
    except Exception as e:
        print(f"Error convirtiendo {doc_file} a DOCX: {e}")
        return None



def extract_text_docx(file_path, umbral):
    """Extrae texto de un archivo DOCX con parada temprana si supera el umbral."""
    try:
        doc = Document(file_path)
        texto_acumulado = []
        total_caracteres = 0

        for para in doc.paragraphs:
            parrafo = para.text
            texto_acumulado.append(parrafo)
            total_caracteres += len(parrafo)
            if total_caracteres > umbral:
                break
        return "\n".join(texto_acumulado)
    except Exception as e:
        print(f"Error extrayendo texto de {file_path} (DOCX): {e}")
        return ""

def process_office(file_path, umbral, digitalizados_dir, escaneados_dir, mllm_dir):
    filename = os.path.basename(file_path).encode("utf-8", errors="replace").decode("utf-8")

    filename_lower = filename.lower()

    # Caso especial: si detecta 'hiring', etc.
    if "hiring" in filename_lower and ("cv" in filename_lower or "resume" in filename_lower or "passport" in filename_lower):
        destino = os.path.join(mllm_dir, filename_lower)
        print(f"{filename} renombrado a {filename_lower} y clasificado como escaneado (contiene 'hiring' etc.)")
        try:
            shutil.move(file_path, destino)
        except Exception as e:
            print(f"Error moviendo {filename}: {e}")
        return

    ext = os.path.splitext(filename_lower)[1]
    texto_extraido = ""

    if ext == ".docx":
        texto_extraido = extract_text_docx(file_path, umbral)
    elif ext == ".doc":
        convert_doc_to_docx(file_path)
    elif ext == ".pptx":
        texto_extraido = extract_text_pptx(file_path, umbral=umbral)
    elif ext == ".ppt":
        shutil.move(file_path, escaneados_dir)
        print("Extraer texto de .ppt si lo deseas con una función similar.")
        return
    else:
        print(f"Formato no soportado: {filename}")
        return

    # Lógica de clasificación
    if len(texto_extraido.strip()) > umbral:
        destino = os.path.join(digitalizados_dir, filename)
        print(f"{filename} clasificado como digitalizado.")
    else:
        destino = os.path.join(escaneados_dir, filename)
        print(f"{filename} clasificado como escaneado (no supera umbral).")

    try:
        if os.path.exists(file_path):
            shutil.move(file_path, destino)
    except Exception as e:
        print(f"Error moviendo {filename}: {e}")

def extract_text_pptx(file_path, umbral):
    """Extrae texto de un archivo PPTX con parada temprana si supera el umbral."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texto_acumulado = []
        total_caracteres = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    contenido = shape.text
                    texto_acumulado.append(contenido)
                    total_caracteres += len(contenido)
                    if total_caracteres > umbral:
                        break
            if total_caracteres > umbral:
                break
        return "\n".join(texto_acumulado)
    except Exception as e:
        print(f"Error extrayendo texto de {file_path} (PPTX): {e}")
        return ""


def filtrado_docx(root_dir):

    digitalizados_dir = os.path.join(root_dir, "digitalizados")
    escaneados_dir = os.path.join(root_dir, "escaneados")
    mllm_dir = os.path.join(root_dir, "to_MLLM")
    umbral = 50

    os.makedirs(digitalizados_dir, exist_ok=True)
    os.makedirs(escaneados_dir, exist_ok=True)
    os.makedirs(mllm_dir, exist_ok=True)

    extensiones = [".doc", ".docx", ".ppt", ".pptx"]


    files = [
        os.path.join(root_dir, f)
        for f in os.listdir(root_dir)
        if os.path.splitext(f.lower())[1] in extensiones
    ]


    with ThreadPoolExecutor(max_workers=10) as executor:
        futures = [executor.submit(
            process_office, file, umbral, digitalizados_dir, escaneados_dir, mllm_dir
        ) for file in files]
        for future in futures:
            future.result()

    return digitalizados_dir, escaneados_dir, mllm_dir



if __name__ == "__main__":
    root_dir = "Dir_documentos"
    digitalizados_dir = os.path.join(root_dir, "digitalizados")
    escaneados_dir = os.path.join(root_dir, "escaneados")
    mllm_dir = os.path.join(root_dir, "to_MLLM")
    umbral = 50

    os.makedirs(digitalizados_dir, exist_ok=True)
    os.makedirs(escaneados_dir, exist_ok=True)
    os.makedirs(mllm_dir, exist_ok=True)

    extensiones = [".doc", ".docx", ".ppt", ".pptx"]

    from concurrent.futures import ThreadPoolExecutor

    files = [
        os.path.join(root_dir, f)
        for f in os.listdir(root_dir)
        if os.path.splitext(f.lower())[1] in extensiones
    ]

    with ThreadPoolExecutor(max_workers=10) as executor:
        futures = [executor.submit(
            process_office, file, umbral, digitalizados_dir, escaneados_dir, mllm_dir
        ) for file in files]
        for future in futures:
            future.result()
